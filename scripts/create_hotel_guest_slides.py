from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUTPUT = ROOT / "Jellin_Hotel_Guest_Use_Case.pptx"

COLORS = {
    "ink": RGBColor(24, 33, 53),
    "muted": RGBColor(90, 102, 124),
    "purple": RGBColor(109, 40, 217),
    "pink": RGBColor(219, 39, 119),
    "soft": RGBColor(245, 243, 255),
    "line": RGBColor(221, 214, 254),
    "white": RGBColor(255, 255, 255),
    "green": RGBColor(22, 163, 74),
}


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(6.8), Inches(1.0))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.bold = True
    run.font.size = Pt(26)
    run.font.color.rgb = COLORS["ink"]

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.72), Inches(1.28), Inches(6.3), Inches(0.5))
        sub_frame = sub.text_frame
        sub_frame.clear()
        p = sub_frame.paragraphs[0]
        run = p.add_run()
        run.text = subtitle
        run.font.name = "Arial"
        run.font.size = Pt(12)
        run.font.color.rgb = COLORS["purple"]


def add_bullets(slide, items, left, top, width, height, font_size=18, color_key="ink"):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.clear()
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = "Arial"
        p.font.size = Pt(font_size)
        p.font.color.rgb = COLORS[color_key]
        p.space_after = Pt(8)
        p.bullet = True
    return box


def add_caption_card(slide, title, body, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS["white"]
    shape.line.color.rgb = COLORS["line"]

    frame = shape.text_frame
    frame.word_wrap = True
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.TOP

    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = "Arial"
    run.font.bold = True
    run.font.size = Pt(18)
    run.font.color.rgb = COLORS["purple"]

    p = frame.add_paragraph()
    p.text = body
    p.font.name = "Arial"
    p.font.size = Pt(13)
    p.font.color.rgb = COLORS["muted"]
    p.space_before = Pt(6)


def add_kicker(slide, text):
    pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(0.22), Inches(2.2), Inches(0.34))
    pill.fill.solid()
    pill.fill.fore_color.rgb = COLORS["soft"]
    pill.line.color.rgb = COLORS["soft"]
    frame = pill.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = COLORS["purple"]


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    hotel = ASSETS / "jellin-hotel.jpeg"
    check_in = ASSETS / "Treat-app-check-in-arrow.jpg"
    loyalty = ASSETS / "Treat-app-loyalty-programs.jpg"
    coupons = ASSETS / "treat-app-coupons-list.jpg"
    qr_coupon = ASSETS / "jellin-qr-coupon-01.jpg"
    cafe = ASSETS / "jellin-everyday-regular-01.jpeg"

    # Slide 1
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COLORS["white"])
    add_kicker(slide, "JELLIN HOSPITALITY USE CASE")
    add_title(slide, "The Hotel Guest Scenario", "A hotel turns one check-in into breakfast, coffee, and neighborhood discovery.")
    slide.shapes.add_picture(str(hotel), Inches(7.7), Inches(0.45), width=Inches(5.0), height=Inches(6.2))
    add_bullets(
        slide,
        [
            "Guest shares only a phone number at reception.",
            "Jellin sends an instant welcome reward by SMS or QR link.",
            "The same loyalty identity extends to partner businesses nearby.",
        ],
        Inches(0.8),
        Inches(2.0),
        Inches(6.2),
        Inches(2.4),
        font_size=20,
    )
    add_caption_card(
        slide,
        "Core promise",
        "No app download for the traveler. Immediate value for the hotel. Measurable spillover to local partners.",
        Inches(0.8),
        Inches(5.05),
        Inches(5.8),
        Inches(1.2),
    )

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COLORS["soft"])
    add_title(slide, "Before Jellin vs With Jellin", "The hospitality problem is not awareness. It is friction and poor follow-through.")
    add_caption_card(
        slide,
        "Before Jellin",
        "A traveler checks into an unfamiliar city, gets paper flyers or generic suggestions, and rarely engages with nearby partners.",
        Inches(0.8),
        Inches(1.7),
        Inches(5.7),
        Inches(2.0),
    )
    add_caption_card(
        slide,
        "With Jellin",
        "The concierge Jells the guest at check-in. The guest receives 5 instant Jells, breakfast value, a free coffee, and nearby partner discounts in seconds.",
        Inches(0.8),
        Inches(4.0),
        Inches(5.7),
        Inches(2.0),
    )
    slide.shapes.add_picture(str(check_in), Inches(7.0), Inches(1.55), width=Inches(5.4), height=Inches(4.6))
    quote = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.55), Inches(5.9), Inches(4.6), Inches(0.7))
    quote.fill.solid()
    quote.fill.fore_color.rgb = COLORS["white"]
    quote.line.color.rgb = COLORS["line"]
    tf = quote.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = '"Wow, that\'s great. Can I use it at other places too?"'
    run.font.name = "Arial"
    run.font.size = Pt(16)
    run.font.italic = True
    run.font.color.rgb = COLORS["pink"]

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COLORS["white"])
    add_title(slide, "Journey Step 1: Reception Check-In", "The concierge activates the guest with almost no training and no signup form.")
    slide.shapes.add_picture(str(check_in), Inches(0.8), Inches(1.5), width=Inches(6.0), height=Inches(4.8))
    add_bullets(
        slide,
        [
            "Concierge asks: 'Your number?'",
            "Operator checks the guest in from the Jellin business flow.",
            "Guest is enrolled under the hotel's loyalty umbrella instantly.",
            "The hotel can assign Jells and instant coupons at the same moment.",
        ],
        Inches(7.2),
        Inches(1.7),
        Inches(5.1),
        Inches(2.8),
        font_size=18,
    )
    add_caption_card(
        slide,
        "Operational benefit",
        "Front desk staff can create a premium first-touch moment without asking the guest to download an app or fill out a long form.",
        Inches(7.2),
        Inches(5.0),
        Inches(5.1),
        Inches(1.3),
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COLORS["soft"])
    add_title(slide, "Journey Step 2: The Guest Receives Value Immediately", "Jellin converts check-in into visible rewards the traveler can actually use.")
    slide.shapes.add_picture(str(loyalty), Inches(0.9), Inches(1.7), width=Inches(3.8), height=Inches(4.9))
    slide.shapes.add_picture(str(coupons), Inches(4.95), Inches(1.7), width=Inches(3.8), height=Inches(4.9))
    add_bullets(
        slide,
        [
            "Breakfast discount at the hotel restaurant.",
            "Free coffee reward as an immediate delight moment.",
            "Partner offers from nearby cafe and boutique under the same ecosystem.",
        ],
        Inches(9.1),
        Inches(2.0),
        Inches(3.4),
        Inches(2.3),
        font_size=17,
    )
    add_caption_card(
        slide,
        "Guest experience",
        "The traveler understands the reward in seconds because the value is concrete, local, and usable right away.",
        Inches(9.0),
        Inches(4.85),
        Inches(3.5),
        Inches(1.3),
    )

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COLORS["white"])
    add_title(slide, "Journey Step 3: Redeem Across the Neighborhood", "One hotel-sponsored identity creates traffic for partner businesses.")
    slide.shapes.add_picture(str(cafe), Inches(0.9), Inches(1.7), width=Inches(5.0), height=Inches(4.6))
    slide.shapes.add_picture(str(qr_coupon), Inches(6.3), Inches(1.7), width=Inches(2.8), height=Inches(4.6))
    add_bullets(
        slide,
        [
            "Guest visits a partner cafe the next morning.",
            "They show the QR reward from SMS, email, or app.",
            "The partner redeems the offer without owning a separate loyalty program.",
            "The hotel measures partner engagement instead of handing out untracked flyers.",
        ],
        Inches(9.4),
        Inches(1.95),
        Inches(3.0),
        Inches(3.5),
        font_size=16,
    )
    add_caption_card(
        slide,
        "Cross-business engine",
        "This is the distinctive Jellin move: the hotel becomes the hub of a local rewards network instead of a standalone property.",
        Inches(6.25),
        Inches(5.4),
        Inches(6.1),
        Inches(1.15),
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_bg(slide, COLORS["soft"])
    add_title(slide, "Why This Use Case Matters", "Hotel guest onboarding becomes a commercial channel, not just a service moment.")
    add_bullets(
        slide,
        [
            "Increase guest satisfaction with immediate, relevant welcome rewards.",
            "Drive breakfast, coffee, and on-site upsell conversion.",
            "Send measurable traffic to neighborhood partners and concession points.",
            "Operate with phone number, email, or QR flows depending on guest preference.",
            "Keep the experience privacy-conscious and low-friction.",
        ],
        Inches(0.9),
        Inches(1.75),
        Inches(6.1),
        Inches(4.3),
        font_size=18,
    )
    panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.4), Inches(1.75), Inches(4.9), Inches(4.7))
    panel.fill.solid()
    panel.fill.fore_color.rgb = COLORS["purple"]
    panel.line.color.rgb = COLORS["purple"]
    tf = panel.text_frame
    tf.clear()
    tf.word_wrap = True
    lines = [
        ("Ideal for", True),
        ("Hotels, resorts, serviced apartments, mall-linked hospitality hubs, and neighborhood partnerships.", False),
        ("Use in pitch", True),
        ("Show how Jellin connects guest experience, partner discovery, and measurable local commerce in one flow.", False),
    ]
    for i, (text, is_header) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = "Arial"
        p.font.size = Pt(18 if is_header else 15)
        p.font.bold = is_header
        p.font.color.rgb = COLORS["white"]
        p.space_after = Pt(10)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build_deck()
    print(f"Created {path}")